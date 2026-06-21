from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
DARK   = RGBColor(0x12, 0x12, 0x1F)   # near-black navy
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x29, 0x7F, 0xFF)   # accent blue
LIGHT  = RGBColor(0xF4, 0xF7, 0xFF)   # off-white bg
GREY   = RGBColor(0x6B, 0x7A, 0x99)   # muted text
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xE7, 0x4C, 0x3C)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=Pt(18), bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bg(slide, color=LIGHT):
    add_rect(slide, 0, 0, W, H, fill=color)

def accent_bar(slide, y=Inches(0.55), color=BLUE):
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.9), fill=color)

def slide_title(slide, text, x=Inches(0.72), y=Inches(0.4), w=Inches(11)):
    add_text(slide, text, x, y, w, Inches(0.7),
             size=Pt(32), bold=True, color=DARK)

def bullet(slide, items, x=Inches(0.72), y_start=Inches(1.55),
           line_h=Inches(0.62), size=Pt(19), color=DARK, dot_color=BLUE):
    for i, item in enumerate(items):
        yy = y_start + i * line_h
        # dot
        add_rect(slide, x, yy + Inches(0.18), Inches(0.1), Inches(0.1), fill=dot_color)
        add_text(slide, item, x + Inches(0.2), yy, Inches(11.8), line_h,
                 size=size, color=color)

def chip(slide, text, x, y, w=Inches(2.4), h=Inches(0.55),
         bg_color=BLUE, txt_color=WHITE, size=Pt(15), bold=False):
    add_rect(slide, x, y, w, h, fill=bg_color)
    add_text(slide, text, x, y + Inches(0.06), w, h,
             size=size, bold=bold, color=txt_color, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. Title ───────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill=DARK)
add_rect(s, 0, Inches(5.8), W, Inches(1.7), fill=BLUE)   # bottom strip
add_text(s, "Invoice Compliance Agent", Inches(1), Inches(1.8), Inches(11.3), Inches(1.4),
         size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Automated grant invoice auditing — end to end",
         Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
         size=Pt(24), color=RGBColor(0xA8, 0xC8, 0xEA), align=PP_ALIGN.CENTER)
add_text(s, "SociAI · Demo Day 2026",
         Inches(1), Inches(6.1), Inches(11.3), Inches(0.6),
         size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)

# ── 2. The Problem ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "The Problem")
bullet(s, [
    "Hundreds of invoices per grant cycle — reviewed manually",
    "29+ compliance rules across vendor, amount, format & document checks",
    "A single missed rule → rejected reimbursement claim",
    "Auditors spend days on work that could take minutes",
])
add_rect(s, Inches(0.72), Inches(5.8), Inches(11.9), Inches(0.8), fill=RGBColor(0xFF, 0xF0, 0xF0))
add_text(s, "Manual review doesn't scale.", Inches(0.9), Inches(5.85), Inches(11), Inches(0.7),
         size=Pt(18), bold=True, color=RED)

# ── 3. The Solution ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "The Solution")

add_text(s, "PDF invoice  →  Invoice Agent  →  Compliance Report",
         Inches(0.72), Inches(1.55), Inches(11.9), Inches(0.7),
         size=Pt(22), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

cols = [
    ("Extract",   "Reads every field\nfrom any layout"),
    ("Evaluate",  "29+ rules checked\nin code & with AI"),
    ("Escalate",  "Flags edge cases\nfor human review"),
    ("Learn",     "Improves with\nevery run"),
]
cx = Inches(0.6)
for label, desc in cols:
    add_rect(s, cx, Inches(2.6), Inches(2.9), Inches(2.8), fill=DARK)
    add_text(s, label, cx, Inches(2.75), Inches(2.9), Inches(0.8),
             size=Pt(22), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(s, desc, cx, Inches(3.6), Inches(2.9), Inches(1.2),
             size=Pt(16), color=WHITE, align=PP_ALIGN.CENTER)
    cx += Inches(3.18)

add_text(s, "Covers: Travel · Local Personnel · HQ Personnel · Equipment · Consumables",
         Inches(0.72), Inches(6.5), Inches(11.9), Inches(0.6),
         size=Pt(15), color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ── 4. Workflow ────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Workflow — Three Phases")

phases = [
    ("1  SCAN",    BLUE,                       "Low-res page map\nClassify document type\nBuild extraction plan"),
    ("2  EXTRACT", RGBColor(0x85, 0x29, 0xFF), "Full-quality render\nOCR + Vision AI per page\nCrop & retry on misses"),
    ("3  VALIDATE",GREEN,                      "Run all 29 rules\nVisual checks (stamps, sigs)\nFlag or finish"),
]
bx = Inches(0.55)
for label, color, desc in phases:
    add_rect(s, bx, Inches(1.8), Inches(3.85), Inches(3.5), fill=color)
    add_text(s, label, bx, Inches(1.9), Inches(3.85), Inches(0.75),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, bx + Inches(0.15), Inches(2.75), Inches(3.5), Inches(2.3),
             size=Pt(17), color=WHITE)
    bx += Inches(4.05)

add_text(s, "Each phase exposes only the tools needed — the agent cannot skip steps.",
         Inches(0.72), Inches(6.4), Inches(11.9), Inches(0.7),
         size=Pt(15), color=GREY, italic=True)

# ── 5. Phase 1 — Scan ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
add_rect(s, 0, 0, Inches(0.25), H, fill=BLUE)
slide_title(s, "Phase 1 — Scan", x=Inches(0.55))
bullet(s, [
    "inspect_file   →   read metadata, check file size",
    "compress_pages   →   render all pages at 48 DPI",
    "inventory_pages   →   vision AI labels each page role",
    "classify_document_type   →   identify invoice type",
    "generate_plan   →   reasoning model writes step-by-step plan",
], x=Inches(0.55), size=Pt(18))
add_text(s, "Fast & cheap before any full-quality work begins.",
         Inches(0.55), Inches(6.35), Inches(11.9), Inches(0.7),
         size=Pt(16), color=GREY, italic=True)

# ── 6. Phase 2 — Extract ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
add_rect(s, 0, 0, Inches(0.25), H, fill=RGBColor(0x85, 0x29, 0xFF))
slide_title(s, "Phase 2 — Extract", x=Inches(0.55))
bullet(s, [
    "convert_pdf_to_images   →   200 DPI render",
    "extract_fields_vision   →   Vision LLM + configured OCR pre-pass",
    "crop_region   →   zoom into missed areas and retry",
    "check_compliance   →   deterministic rule pass (no LLM)",
    "check_compliance_visual   →   stamps, seals, signatures",
], x=Inches(0.55), size=Pt(18))
add_text(s, "Confidence-aware merging — never overwrites a higher-confidence value.",
         Inches(0.55), Inches(6.35), Inches(11.9), Inches(0.7),
         size=Pt(16), color=GREY, italic=True)

# ── 7. Phase 3 — Validate ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
add_rect(s, 0, 0, Inches(0.25), H, fill=GREEN)
slide_title(s, "Phase 3 — Validate & Finish", x=Inches(0.55))
bullet(s, [
    "Re-run compliance after any new extraction",
    "flag_for_human_review   →   unresolvable cases escalated",
    "write_learning   →   persist insights for future runs",
    "finish   →   compliance_passed · human_review_needed · max_retries",
], x=Inches(0.55), size=Pt(18))

add_rect(s, Inches(0.55), Inches(5.3), Inches(11.9), Inches(1.7), fill=DARK)
add_text(s, "Output per invoice", Inches(0.8), Inches(5.35), Inches(5), Inches(0.5),
         size=Pt(15), bold=True, color=BLUE)
add_text(s,
    "results.csv  ·  compliance.csv  ·  page images  ·  full agent log (JSONL)",
    Inches(0.8), Inches(5.85), Inches(11.5), Inches(0.9),
    size=Pt(15), color=WHITE)

# ── 8. The Learning System ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Built-in Learning")

add_text(s, "After every run (--learn mode):", Inches(0.72), Inches(1.55),
         Inches(11.9), Inches(0.5), size=Pt(18), color=GREY)
bullet(s, [
    "Reflection loop compares output against ground truth",
    "Writes targeted learnings per invoice type to learnings.md",
    "On the next run, relevant learnings are injected into the system prompt",
], y_start=Inches(2.2), size=Pt(19))

add_text(s, "Learning categories", Inches(0.72), Inches(4.0), Inches(4), Inches(0.5),
         size=Pt(16), bold=True, color=DARK)
cats = ["extraction patterns", "vision model quirks", "common failures",
        "compliance edge cases", "general approaches"]
cy = Inches(4.6)
cx = Inches(0.72)
for cat in cats:
    chip(s, cat, cx, cy, w=Inches(2.55), h=Inches(0.48),
         bg_color=RGBColor(0xE8, 0xF0, 0xFF), txt_color=BLUE, size=Pt(13))
    cx += Inches(2.7)
    if cx > Inches(11):
        cx = Inches(0.72)
        cy += Inches(0.6)

add_text(s, "The agent compounds — each invoice makes the next one more accurate.",
         Inches(0.72), Inches(6.4), Inches(11.9), Inches(0.7),
         size=Pt(15), color=GREY, italic=True)

# ── 9. Technical Architecture ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Under the Hood")

# Left column — stack diagram
boxes = [
    ("main.py", BLUE),
    ("InvoiceAgent", RGBColor(0x29, 0x5F, 0xBF)),
    ("AgentState  ·  ToolRegistry  ·  LLMProvider", RGBColor(0x1A, 0x3A, 0x7A)),
    ("agent_turn()  →  prompt  →  JSON action  →  tool dispatch", DARK),
]
by = Inches(1.7)
for label, color in boxes:
    add_rect(s, Inches(0.6), by, Inches(7.2), Inches(0.72), fill=color)
    add_text(s, label, Inches(0.75), by + Inches(0.1), Inches(7), Inches(0.55),
             size=Pt(15), color=WHITE, bold=True)
    by += Inches(0.85)

# Right column — model split
add_text(s, "Dual model split", Inches(8.4), Inches(1.55), Inches(4.5), Inches(0.5),
         size=Pt(18), bold=True, color=DARK)
rows = [
    ("Reasoning", "gemini-2.5-flash", "Tool selection · Planning"),
    ("Vision",    "gemini-2.5-pro",   "Extraction · Compliance"),
]
ry = Inches(2.2)
for role, model, use in rows:
    add_rect(s, Inches(8.4), ry, Inches(4.5), Inches(1.05), fill=DARK)
    add_text(s, role, Inches(8.55), ry + Inches(0.05), Inches(2), Inches(0.45),
             size=Pt(14), bold=True, color=BLUE)
    add_text(s, model, Inches(8.55), ry + Inches(0.45), Inches(4.2), Inches(0.4),
             size=Pt(13), color=WHITE)
    add_text(s, use, Inches(8.55), ry + Inches(0.72), Inches(4.2), Inches(0.3),
             size=Pt(11), color=GREY)
    ry += Inches(1.2)

add_text(s, "Also runs fully local via Ollama — no data leaves the machine.",
         Inches(8.4), Inches(4.8), Inches(4.5), Inches(0.6),
         size=Pt(13), color=GREY, italic=True)

# guard boxes
add_text(s, "Safety guards", Inches(0.72), Inches(5.15), Inches(7), Inches(0.45),
         size=Pt(16), bold=True, color=DARK)
guards = ["Duplicate action detection", "Consecutive failure limit",
          "JSON repair-retry", "Per-run token cap"]
gx = Inches(0.72)
for g in guards:
    chip(s, g, gx, Inches(5.7), w=Inches(3.05), h=Inches(0.48),
         bg_color=RGBColor(0xE8, 0xF0, 0xFF), txt_color=BLUE, size=Pt(12))
    gx += Inches(3.2)

# ── 10. Config-Driven ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Zero Code to Extend")

add_text(s, "All invoice types, fields, and rules live in CSV files.",
         Inches(0.72), Inches(1.55), Inches(11.9), Inches(0.5),
         size=Pt(19), color=GREY)

files = [
    ("invoice_types.csv",    "5 types with LLM context hints"),
    ("extraction_fields.csv","Per-type fields, page regions, hints"),
    ("compliance_rules.csv", "29+ rules: required, regex, range,\ncross-field, conditional, visual"),
    ("phase_tools.yaml",     "Which tools are visible in each phase"),
]
fx = Inches(0.55)
for fname, desc in files:
    add_rect(s, fx, Inches(2.4), Inches(2.95), Inches(2.8), fill=DARK)
    add_text(s, fname, fx + Inches(0.1), Inches(2.5), Inches(2.75), Inches(0.55),
             size=Pt(14), bold=True, color=BLUE)
    add_text(s, desc, fx + Inches(0.1), Inches(3.1), Inches(2.75), Inches(1.8),
             size=Pt(14), color=WHITE)
    fx += Inches(3.18)

add_text(s, "Add a new invoice type: one row in each CSV. No Python required.",
         Inches(0.72), Inches(6.4), Inches(11.9), Inches(0.7),
         size=Pt(16), color=GREY, italic=True)

# ── 11. Demo ──────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
slide_title(s, "Demo")

bullet(s, [
    "Input: scanned PDF travel invoice (VIAJES type)",
    "Live: agent phases through Scan → Extract → Validate",
    "Output: extracted fields with confidence + compliance report",
    "Bonus: agent reasoning trace — why each tool was chosen",
], size=Pt(20), line_h=Inches(0.85), y_start=Inches(1.7))

# ── 12. Thank You ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, W, H, fill=DARK)
add_rect(s, 0, Inches(5.8), W, Inches(1.7), fill=BLUE)
add_text(s, "Thank You", Inches(1), Inches(2.0), Inches(11.3), Inches(1.4),
         size=Pt(60), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Questions?",
         Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
         size=Pt(26), color=RGBColor(0xA8, 0xC8, 0xEA), align=PP_ALIGN.CENTER)
add_text(s, "Invoice Compliance Agent — SociAI 2026",
         Inches(1), Inches(6.1), Inches(11.3), Inches(0.6),
         size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/juliankraus/Coding/SociAI/invoice-agent/docs/invoice_agent_demo.pptx"
prs.save(out)
print(f"Saved: {out}")
